"""Chrome extension API — manages audit sessions, observations, and scoring.

The Chrome extension records human tester interactions with AI tools and
uploads structured observation data for LLM-based scoring.

Auth: API key with 'agent:write' scope (X-API-Key header).
"""

import base64
import hashlib
import json
import logging
import os
import re
import uuid
from datetime import datetime, timezone
from pathlib import Path

from fastapi import APIRouter, BackgroundTasks, Depends, File, HTTPException, UploadFile
from sqlalchemy import text
from sqlalchemy.ext.asyncio import AsyncSession

from ..deps import get_db
from .agent import require_agent_key
from ...db.models.user import User
from ...schemas.extension import (
    ExtensionSessionCreate,
    ExtensionSessionResponse,
    ObservationResponse,
    ObservationUpload,
    SessionProgressResponse,
    TestCaseOut,
    ToolListItem,
)

logger = logging.getLogger(__name__)

router = APIRouter()

# UUID format regex for session_id validation (prevents path traversal)
_UUID_RE = re.compile(r'^[a-f0-9]{8}-[a-f0-9]{4}-[a-f0-9]{4}-[a-f0-9]{4}-[a-f0-9]{12}$')


def _validate_session_id(session_id: str) -> None:
    """Validate session_id is a proper UUID to prevent path traversal."""
    if not _UUID_RE.match(session_id):
        raise HTTPException(400, "Invalid session ID format")

# Screenshots storage base path (persistent volume in production)
from ...config import settings as _ext_settings
_SCREENSHOTS_DIR = Path(_ext_settings.screenshots_dir)


# ---------------------------------------------------------------------------
# POST /sessions — Create extension audit session
# ---------------------------------------------------------------------------

@router.post("/sessions", response_model=ExtensionSessionResponse)
async def create_extension_session(
    body: ExtensionSessionCreate,
    db: AsyncSession = Depends(get_db),
    user: User = Depends(require_agent_key),
):
    """Create a new audit session for the Chrome extension."""
    try:
        return await _create_session_impl(body, db, user)
    except HTTPException:
        raise
    except Exception as e:
        logger.exception("Session creation failed: %s", e)
        raise HTTPException(500, f"セッション作成に失敗: {type(e).__name__}: {e}")


async def _create_session_impl(body, db, user):
    session_id = str(uuid.uuid4())
    now = datetime.utcnow()

    # Generate session code
    hash_input = f"{now.isoformat()}-{session_id}"
    short_hash = hashlib.sha256(hash_input.encode()).hexdigest()[:8].upper()
    date_str = now.strftime("%Y%m%d")
    session_code = f"AX-{date_str}-{short_hash}"

    # Verify tool exists
    result = await db.execute(
        text("SELECT id, name FROM tools WHERE id = :tid"),
        {"tid": body.tool_id},
    )
    tool_row = result.fetchone()
    if not tool_row:
        raise HTTPException(404, f"ツールが見つかりません: {body.tool_id}")

    # Generate test cases in memory first (don't insert yet — FK requires session)
    test_cases_out: list[TestCaseOut] = []
    cases = []
    total_planned = 0

    if body.recording_mode == "protocol":
        try:
            from aixis_agent.patterns.generator import generate_all

            config_dir = Path(__file__).resolve().parents[4] / "config"
            patterns_dir = config_dir / "patterns"

            # Resolve categories from profile if not explicitly provided
            categories = body.categories
            if not categories and body.profile_id:
                try:
                    from aixis_agent.profiles.registry import (
                        get_categories_for_profile,
                        get_profile,
                    )
                    profiles_dir = config_dir / "profiles"
                    profile = get_profile(body.profile_id, profiles_dir)
                    if profile:
                        categories = get_categories_for_profile(profile)
                except Exception:
                    pass

            cases = generate_all(patterns_dir, categories)

            # Sort by priority for optimal coverage
            from aixis_agent.orchestrator.pipeline import sort_by_priority
            cases = sort_by_priority(cases)

            # Cap test cases for manual testing (human testers can't do 500+)
            max_cases = body.max_cases or 30
            if len(cases) > max_cases:
                cases = cases[:max_cases]

            total_planned = len(cases)

        except Exception as e:
            logger.exception("Test case generation failed: %s", e)
            raise HTTPException(500, f"テストケース生成に失敗しました: {type(e).__name__}: {e}")

    if not cases:
        raise HTTPException(500, "テストケースが0件です。パターンファイルの設定を確認してください。")

    # 1. Create session FIRST (test cases have FK to audit_sessions)
    await db.execute(text("""
        INSERT INTO audit_sessions
        (id, session_code, tool_id, profile_id, status, initiated_by,
         created_at, executor_type, total_planned, total_executed)
        VALUES (:id, :code, :tool_id, :profile_id, :status, :initiated_by,
                :now, :executor_type, :total_planned, 0)
    """), {
        "id": session_id,
        "code": session_code,
        "tool_id": body.tool_id,
        "profile_id": body.profile_id or "",
        "status": "running",
        "initiated_by": user.id,
        "now": now,
        "executor_type": "extension",
        "total_planned": total_planned,
    })

    # 2. Insert test cases AFTER session exists
    for case in cases:
        cat_val = case.category.value if hasattr(case.category, "value") else str(case.category)
        try:
            await db.execute(text("""
                INSERT INTO db_test_cases
                (id, session_id, category, prompt, metadata_json,
                 expected_behaviors, failure_indicators, tags)
                VALUES (:id, :session_id, :category, :prompt, :metadata,
                        :expected, :failures, :tags)
                ON CONFLICT (id) DO NOTHING
            """), {
                "id": case.id,
                "session_id": session_id,
                "category": cat_val,
                "prompt": case.prompt,
                "metadata": json.dumps(case.metadata, ensure_ascii=False),
                "expected": json.dumps(case.expected_behaviors, ensure_ascii=False),
                "failures": json.dumps(case.failure_indicators, ensure_ascii=False),
                "tags": json.dumps(case.tags, ensure_ascii=False),
            })
        except Exception as e:
            logger.warning("Failed to insert test case %s: %s", case.id, e)
            continue

        test_cases_out.append(TestCaseOut(
            id=case.id,
            category=cat_val,
            prompt=case.prompt,
            expected_behaviors=case.expected_behaviors,
            failure_indicators=case.failure_indicators,
            tags=case.tags,
            metadata=case.metadata,
        ))

    await db.commit()

    logger.info(
        "Extension session created: %s (tool=%s, mode=%s, cases=%d, user=%s)",
        session_code, body.tool_id, body.recording_mode, total_planned, user.email,
    )
    return ExtensionSessionResponse(
        session_id=session_id,
        session_code=session_code,
        tool_id=body.tool_id,
        status="running",
        recording_mode=body.recording_mode,
        test_cases=test_cases_out,
    )


# ---------------------------------------------------------------------------
# GET /sessions/{id}/test-cases — Get test cases for session
# ---------------------------------------------------------------------------

@router.get("/sessions/{session_id}/test-cases", response_model=list[TestCaseOut])
async def get_session_test_cases(
    session_id: str,
    db: AsyncSession = Depends(get_db),
    user: User = Depends(require_agent_key),
):
    """Return test cases for a session."""
    _validate_session_id(session_id)
    result = await db.execute(
        text("""
            SELECT id, category, prompt, metadata_json,
                   expected_behaviors, failure_indicators, tags
            FROM db_test_cases
            WHERE session_id = :sid
            ORDER BY id
        """),
        {"sid": session_id},
    )
    rows = result.fetchall()
    if not rows:
        raise HTTPException(404, "テストケースが見つかりません")

    cases = []
    for row in rows:
        cases.append(TestCaseOut(
            id=row[0],
            category=row[1],
            prompt=row[2],
            metadata=row[3] if isinstance(row[3], dict) else (json.loads(row[3]) if row[3] else {}),
            expected_behaviors=row[4] if isinstance(row[4], list) else (json.loads(row[4]) if row[4] else []),
            failure_indicators=row[5] if isinstance(row[5], list) else (json.loads(row[5]) if row[5] else []),
            tags=row[6] if isinstance(row[6], list) else (json.loads(row[6]) if row[6] else []),
        ))
    return cases


# ---------------------------------------------------------------------------
# POST /sessions/{id}/observations — Upload observation
# ---------------------------------------------------------------------------

@router.post("/sessions/{session_id}/observations", response_model=ObservationResponse)
async def upload_observation(
    session_id: str,
    body: ObservationUpload,
    db: AsyncSession = Depends(get_db),
    user: User = Depends(require_agent_key),
):
    """Upload a single observation (input/output pair) from the Chrome extension."""
    _validate_session_id(session_id)

    # Limit payload size
    import sys
    body_size = sys.getsizeof(body.model_dump_json())
    if body_size > 10 * 1024 * 1024:  # 10MB limit
        raise HTTPException(400, "ペイロードが大きすぎます（最大10MB）")

    # Verify session exists and belongs to user
    result = await db.execute(
        text("SELECT id, tool_id, status, total_executed FROM audit_sessions WHERE id = :sid AND initiated_by = :uid"),
        {"sid": session_id, "uid": user.id},
    )
    session_row = result.fetchone()
    # Admin/analyst fallback — can access any session
    if not session_row and user.role in ('admin', 'analyst'):
        result = await db.execute(
            text("SELECT id, tool_id, status, total_executed FROM audit_sessions WHERE id = :sid"),
            {"sid": session_id},
        )
        session_row = result.fetchone()
    if not session_row:
        raise HTTPException(404, f"セッションが見つかりません: {session_id}")

    if session_row[2] not in ("running", "pending"):
        raise HTTPException(400, f"セッションは現在 {session_row[2]} 状態です。観察データを追加できません。")

    now = datetime.utcnow()

    # Get actual observation count (including manual screenshots) for sequence numbering
    count_result = await db.execute(
        text("SELECT COUNT(*) FROM db_test_results WHERE session_id = :sid"),
        {"sid": session_id},
    )
    observation_count = count_result.scalar() or 0
    sequence_number = observation_count + 1

    # Handle screenshot
    screenshot_path = None
    if body.screenshot_base64:
        try:
            img_data = base64.b64decode(body.screenshot_base64)
            session_dir = _SCREENSHOTS_DIR / session_id
            session_dir.mkdir(parents=True, exist_ok=True)
            # Detect image format from magic bytes to use correct extension
            img_ext = "png"
            if img_data[:3] == b'\xff\xd8\xff':
                img_ext = "jpg"
            img_path = session_dir / f"{sequence_number:04d}.{img_ext}"
            img_path.write_bytes(img_data)
            # Path relative to screenshots mount; accessible via /screenshots/...
            screenshot_path = f"/screenshots/{session_id}/{sequence_number:04d}.{img_ext}"
        except Exception as e:
            logger.warning("Screenshot save failed: %s", e)

    # Determine observation type: screenshot evidence vs test progression
    capture_type = (body.metadata or {}).get("capture_type", "")
    is_screenshot_evidence = capture_type in (
        "manual_screenshot", "full_screenshot", "partial_screenshot"
    )
    is_manual_screenshot = is_screenshot_evidence

    test_case_id = body.test_case_id
    category = "freeform"

    if is_screenshot_evidence:
        # Screenshots: store as supplementary evidence linked to the current test
        # Use the test_case_id if provided, otherwise create a synthetic one
        if not test_case_id:
            test_case_id = f"screenshot-{session_id[:8]}-{sequence_number:04d}"
        category = "screenshot_evidence"
        # Don't create a synthetic test case — just store the screenshot as evidence
        # The test_case_id links it to the real test case if provided
    elif test_case_id:
        # Protocol mode: look up the actual category from the test case
        tc_result = await db.execute(
            text("SELECT category FROM db_test_cases WHERE id = :tid AND session_id = :sid"),
            {"tid": test_case_id, "sid": session_id},
        )
        tc_row = tc_result.fetchone()
        if tc_row:
            category = tc_row[0]
        else:
            # test_case_id provided but not found in DB — store as-is with protocol category
            logger.warning("Test case %s not found in DB for session %s, creating inline", test_case_id, session_id)
            category = body.metadata.get("category", "protocol") if body.metadata else "protocol"
            await db.execute(text("""
                INSERT INTO db_test_cases
                (id, session_id, category, prompt, metadata_json,
                 expected_behaviors, failure_indicators, tags)
                VALUES (:id, :session_id, :category, :prompt, :metadata,
                        :expected, :failures, :tags)
                ON CONFLICT (id) DO NOTHING
            """), {
                "id": test_case_id,
                "session_id": session_id,
                "category": category,
                "prompt": body.prompt_text or "",
                "metadata": "{}",
                "expected": "[]",
                "failures": "[]",
                "tags": '["protocol"]',
            })
    else:
        # Freeform mode: create synthetic test case
        test_case_id = f"freeform-{session_id[:8]}-{sequence_number:04d}"
        await db.execute(text("""
            INSERT INTO db_test_cases
            (id, session_id, category, prompt, metadata_json,
             expected_behaviors, failure_indicators, tags)
            VALUES (:id, :session_id, :category, :prompt, :metadata,
                    :expected, :failures, :tags)
            ON CONFLICT (id) DO NOTHING
        """), {
            "id": test_case_id,
            "session_id": session_id,
            "category": "freeform",
            "prompt": body.prompt_text or "",
            "metadata": "{}",
            "expected": "[]",
            "failures": "[]",
            "tags": '["freeform"]',
        })

    # Store as DBTestResult
    await db.execute(text("""
        INSERT INTO db_test_results
        (session_id, test_case_id, category, prompt_sent, response_raw,
         response_time_ms, error, screenshot_path, page_url, executed_at, metadata_json)
        VALUES (:session_id, :test_case_id, :category, :prompt, :response,
                :time_ms, :error, :screenshot, :page_url, :executed_at, :metadata)
    """), {
        "session_id": session_id,
        "test_case_id": test_case_id,
        "category": category,
        "prompt": body.prompt_text,
        "response": body.response_text,
        "time_ms": body.response_time_ms,
        "error": None,
        "screenshot": screenshot_path,
        "page_url": body.page_url,
        "executed_at": now,
        "metadata": json.dumps(body.metadata, ensure_ascii=False) if body.metadata else "{}",
    })

    # Use sequence number as observation ID
    obs_id = sequence_number

    # When a test_completion observation arrives with a timer value,
    # backfill response_time_ms onto any earlier screenshot evidence rows
    # for the same test_case_id (they were uploaded with time_ms=0).
    obs_type = (body.metadata or {}).get("type", "")
    if obs_type == "test_completion" and body.response_time_ms and body.response_time_ms > 0:
        await db.execute(text("""
            UPDATE db_test_results
            SET response_time_ms = :time_ms
            WHERE session_id = :sid
              AND test_case_id = :tid
              AND (response_time_ms IS NULL OR response_time_ms = 0)
              AND category = 'screenshot_evidence'
        """), {
            "time_ms": body.response_time_ms,
            "sid": session_id,
            "tid": test_case_id,
        })

    # Update session progress — only count actual test progressions, not screenshots
    if not is_screenshot_evidence:
        await db.execute(text("""
            UPDATE audit_sessions
            SET total_executed = COALESCE(total_executed, 0) + 1,
                started_at = COALESCE(started_at, :now)
            WHERE id = :sid
        """), {
            "now": now,
            "sid": session_id,
        })
    else:
        # Still update started_at if needed
        await db.execute(text("""
            UPDATE audit_sessions
            SET started_at = COALESCE(started_at, :now)
            WHERE id = :sid
        """), {
            "now": now,
            "sid": session_id,
        })

    await db.commit()

    return ObservationResponse(
        observation_id=obs_id,
        sequence_number=sequence_number,
    )


# ---------------------------------------------------------------------------
# POST /sessions/{id}/upload-file — Upload PPTX/PDF artifact for analysis
# ---------------------------------------------------------------------------

@router.post("/sessions/{session_id}/upload-file")
async def upload_file(
    session_id: str,
    file: UploadFile = File(...),
    db: AsyncSession = Depends(get_db),
    user: User = Depends(require_agent_key),
):
    """Upload a PPTX or PDF file and extract text content for LLM scoring."""
    _validate_session_id(session_id)
    # Verify session exists
    result = await db.execute(
        text("SELECT id, status FROM audit_sessions WHERE id = :sid"),
        {"sid": session_id},
    )
    session_row = result.fetchone()
    if not session_row:
        raise HTTPException(404, f"セッションが見つかりません: {session_id}")

    # Validate file type — use only the LAST extension after stripping path
    raw_filename = file.filename or "upload"
    safe_name = os.path.basename(raw_filename)
    ext = safe_name.rsplit(".", 1)[-1].lower() if "." in safe_name else ""
    if ext not in ("pptx", "pdf"):
        raise HTTPException(400, "PPTX または PDF ファイルのみアップロード可能です")

    # Read and validate file size (50MB limit)
    content = await file.read()

    # Validate file magic bytes
    if ext == "pdf" and not content[:4] == b'%PDF':
        raise HTTPException(400, "ファイルの内容がPDF形式ではありません")
    if ext == "pptx" and not content[:2] == b'PK':
        raise HTTPException(400, "ファイルの内容がPPTX形式ではありません")
    if len(content) > 50 * 1024 * 1024:
        raise HTTPException(400, "ファイルサイズが大きすぎます（最大50MB）")

    # Sanitize safe_name — strip path separators and only allow safe characters
    safe_filename = re.sub(r'[^a-zA-Z0-9._-]', '_', safe_name)
    if not safe_filename or safe_filename.startswith('.'):
        safe_filename = f"upload_{uuid.uuid4().hex[:8]}.{ext}"

    # Save file
    upload_dir = Path(__file__).resolve().parents[2] / "static" / "uploads" / "extension" / session_id
    upload_dir.mkdir(parents=True, exist_ok=True)
    file_path = upload_dir / safe_filename
    file_path.write_bytes(content)

    # Extract text
    extracted_text = ""
    try:
        if ext == "pptx":
            from pptx import Presentation
            from io import BytesIO
            prs = Presentation(BytesIO(content))
            slide_texts = []
            for i, slide in enumerate(prs.slides, 1):
                texts = []
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            text_content = para.text.strip()
                            if text_content:
                                texts.append(text_content)
                if texts:
                    slide_texts.append(f"--- スライド {i} ---\n" + "\n".join(texts))
            extracted_text = "\n\n".join(slide_texts)

        elif ext == "pdf":
            from pypdf import PdfReader
            from io import BytesIO
            reader = PdfReader(BytesIO(content))
            page_texts = []
            for i, page in enumerate(reader.pages, 1):
                page_text = page.extract_text()
                if page_text and page_text.strip():
                    page_texts.append(f"--- ページ {i} ---\n{page_text.strip()}")
            extracted_text = "\n\n".join(page_texts)

    except Exception as e:
        logger.warning("File text extraction failed for %s: %s", safe_name, e)
        extracted_text = f"(テキスト抽出に失敗しました: {e})"

    # Store extracted text as a special observation
    now = datetime.utcnow()
    file_case_id = f"file-{session_id[:8]}-{safe_filename}"
    await db.execute(text("""
        INSERT INTO db_test_cases
        (id, session_id, category, prompt, metadata_json,
         expected_behaviors, failure_indicators, tags)
        VALUES (:id, :session_id, :category, :prompt, :metadata,
                :expected, :failures, :tags)
        ON CONFLICT (id) DO NOTHING
    """), {
        "id": file_case_id,
        "session_id": session_id,
        "category": "artifact_upload",
        "prompt": f"アップロードされた成果物: {safe_name}",
        "metadata": json.dumps({"file_type": ext, "filename": safe_name}, ensure_ascii=False),
        "expected": "[]",
        "failures": "[]",
        "tags": json.dumps(["artifact", ext], ensure_ascii=False),
    })

    await db.execute(text("""
        INSERT INTO db_test_results
        (session_id, test_case_id, category, prompt_sent, response_raw,
         response_time_ms, error, screenshot_path, page_url, executed_at, metadata_json)
        VALUES (:session_id, :test_case_id, :category, :prompt, :response,
                0, NULL, NULL, NULL, :now, :metadata)
    """), {
        "session_id": session_id,
        "test_case_id": file_case_id,
        "category": "artifact_upload",
        "prompt": f"成果物ファイル: {safe_name}",
        "response": extracted_text[:50000],  # Limit to 50K chars
        "now": now,
        "metadata": json.dumps({"file_path": str(file_path), "file_size": len(content)}, ensure_ascii=False),
    })

    await db.commit()

    logger.info("File uploaded for session %s: %s (%d bytes, %d chars extracted)",
                session_id, safe_name, len(content), len(extracted_text))

    return {
        "file_id": file_case_id,
        "filename": safe_name,
        "file_type": ext,
        "file_size": len(content),
        "extracted_text_length": len(extracted_text),
        "extracted_text_preview": extracted_text[:500] if extracted_text else "",
    }


# ---------------------------------------------------------------------------
# POST /sessions/{id}/advance — Advance test progress (no observation created)
# ---------------------------------------------------------------------------

@router.post("/sessions/{session_id}/advance")
async def advance_test_progress(
    session_id: str,
    body: dict,
    db: AsyncSession = Depends(get_db),
    user: User = Depends(require_agent_key),
):
    """Advance test progress counter without creating an observation row."""
    _validate_session_id(session_id)

    # Verify session exists, belongs to user, and is in a valid state
    result = await db.execute(
        text("SELECT id, status FROM audit_sessions WHERE id = :sid AND initiated_by = :uid"),
        {"sid": session_id, "uid": user.id},
    )
    session_row = result.fetchone()
    if not session_row and user.role in ('admin', 'analyst'):
        result = await db.execute(
            text("SELECT id, status FROM audit_sessions WHERE id = :sid"),
            {"sid": session_id},
        )
        session_row = result.fetchone()
    if not session_row:
        raise HTTPException(404, f"セッションが見つかりません: {session_id}")
    if session_row[1] not in ("running", "pending"):
        raise HTTPException(400, f"セッションは現在 {session_row[1]} 状態です。進捗を更新できません。")

    # Store response_time_ms for the test case if provided
    response_time_ms = body.get("response_time_ms", 0)
    test_case_id = body.get("test_case_id")
    test_index = body.get("test_index")

    now = datetime.utcnow()

    # Backfill response_time_ms onto db_test_results for this test_case_id
    if response_time_ms and response_time_ms > 0 and test_case_id:
        await db.execute(text("""
            UPDATE db_test_results
            SET response_time_ms = :time_ms
            WHERE session_id = :sid
              AND test_case_id = :tid
              AND (response_time_ms IS NULL OR response_time_ms = 0)
        """), {
            "time_ms": response_time_ms,
            "sid": session_id,
            "tid": test_case_id,
        })

    await db.execute(text("""
        UPDATE audit_sessions
        SET total_executed = COALESCE(total_executed, 0) + 1,
            started_at = COALESCE(started_at, :now)
        WHERE id = :sid
    """), {"now": now, "sid": session_id})
    await db.commit()
    return {"ok": True}


# ---------------------------------------------------------------------------
# POST /sessions/{id}/complete — Mark session complete, trigger scoring
# ---------------------------------------------------------------------------

@router.post("/sessions/{session_id}/complete")
async def complete_session(
    session_id: str,
    background_tasks: BackgroundTasks,
    db: AsyncSession = Depends(get_db),
    user: User = Depends(require_agent_key),
):
    """Mark session as complete and trigger LLM scoring."""
    _validate_session_id(session_id)
    # Verify session exists and belongs to user
    result = await db.execute(
        text("SELECT id, tool_id, status, total_planned, total_executed FROM audit_sessions WHERE id = :sid AND initiated_by = :uid"),
        {"sid": session_id, "uid": user.id},
    )
    session_row = result.fetchone()
    # Admin/analyst fallback — can access any session
    if not session_row and user.role in ('admin', 'analyst'):
        result = await db.execute(
            text("SELECT id, tool_id, status, total_planned, total_executed FROM audit_sessions WHERE id = :sid"),
            {"sid": session_id},
        )
        session_row = result.fetchone()
    if not session_row:
        raise HTTPException(404, f"セッションが見つかりません: {session_id}")

    if session_row[2] not in ("running", "pending"):
        raise HTTPException(400, f"セッションは既に {session_row[2]} 状態です")

    tool_id = session_row[1]
    total_planned = session_row[3] or 0
    total_executed = session_row[4] or 0
    completeness = int(total_executed / total_planned * 100) if total_planned > 0 else 0

    now = datetime.utcnow()

    # Update session to "scoring" status
    await db.execute(text("""
        UPDATE audit_sessions
        SET status = 'scoring',
            completed_at = :now,
            completeness_ratio = :completeness
        WHERE id = :sid
    """), {
        "now": now,
        "completeness": completeness,
        "sid": session_id,
    })
    await db.commit()

    # Schedule LLM scoring via FastAPI BackgroundTasks (reliable execution)
    background_tasks.add_task(_run_llm_scoring_background, session_id, tool_id)

    logger.info("Session %s marked complete, scoring scheduled (observations=%d)", session_id, total_executed)
    return {
        "status": "scoring",
        "session_id": session_id,
        "total_executed": total_executed,
        "completeness_ratio": completeness,
        "message": "LLMスコアリングをバックグラウンドで開始しました",
    }


async def _run_llm_scoring_background(session_id: str, tool_id: str) -> None:
    """Background task for LLM scoring. Runs after response is sent."""
    from ...db.base import async_session

    logger.info("=== LLM scoring START for session %s (tool %s) ===", session_id, tool_id)

    try:
        from ...services.audit_runner import register_job, update_job, cleanup_job
        register_job(session_id, phase="llm_scoring", tool_id=tool_id)
    except Exception:
        pass

    try:
        from ...services.llm_scorer import LLMScorer

        async with async_session() as db:
            scorer = LLMScorer()
            await scorer.score_session(session_id, tool_id, db)

            # Set to completed after scoring
            await db.execute(
                text("UPDATE audit_sessions SET status = 'completed' WHERE id = :sid"),
                {"sid": session_id},
            )
            await db.commit()

        logger.info("=== LLM scoring COMPLETED for session %s ===", session_id)

        try:
            from ...services.audit_runner import update_job, cleanup_job
            update_job(session_id, status="completed", phase="done")
            cleanup_job(session_id)
        except Exception:
            pass

    except Exception as e:
        logger.exception("=== LLM scoring CRASHED for session %s: %s ===", session_id, e)

        # Update session status to failed
        try:
            async with async_session() as db:
                await db.execute(text("""
                    UPDATE audit_sessions
                    SET status = 'failed', error_message = :error
                    WHERE id = :sid
                """), {"error": str(e)[:2000], "sid": session_id})
                await db.commit()
        except Exception as e2:
            logger.error("Failed to update session status for %s: %s", session_id, e2)

        try:
            from ...services.audit_runner import update_job, cleanup_job
            update_job(session_id, status="failed", error=str(e))
            cleanup_job(session_id)
        except Exception:
            pass


# ---------------------------------------------------------------------------
# GET /sessions/{id}/progress — Get session progress
# ---------------------------------------------------------------------------

@router.get("/sessions/{session_id}/progress", response_model=SessionProgressResponse)
async def get_session_progress(
    session_id: str,
    db: AsyncSession = Depends(get_db),
    user: User = Depends(require_agent_key),
):
    """Return current session progress."""
    _validate_session_id(session_id)
    # Verify session exists and belongs to user
    result = await db.execute(
        text("""
            SELECT id, session_code, status, total_planned, total_executed,
                   completeness_ratio, executor_type
            FROM audit_sessions WHERE id = :sid AND initiated_by = :uid
        """),
        {"sid": session_id, "uid": user.id},
    )
    row = result.fetchone()
    # Admin/analyst fallback — can access any session
    if not row and user.role in ('admin', 'analyst'):
        result = await db.execute(
            text("""
                SELECT id, session_code, status, total_planned, total_executed,
                       completeness_ratio, executor_type
                FROM audit_sessions WHERE id = :sid
            """),
            {"sid": session_id},
        )
        row = result.fetchone()
    if not row:
        raise HTTPException(404, f"セッションが見つかりません: {session_id}")

    return SessionProgressResponse(
        session_id=row[0],
        session_code=row[1],
        status=row[2],
        total_planned=row[3] or 0,
        total_executed=row[4] or 0,
        completeness_ratio=row[5] or 0,
        # Inferred from total_planned: protocol mode has test cases, freeform does not
        recording_mode="protocol" if (row[3] or 0) > 0 else "freeform",
    )


# ---------------------------------------------------------------------------
# GET /tools — List available tools for extension UI
# ---------------------------------------------------------------------------

@router.get("/tools", response_model=list[ToolListItem])
async def list_tools_for_extension(
    db: AsyncSession = Depends(get_db),
    user: User = Depends(require_agent_key),
):
    """Return tools available for auditing."""
    try:
        result = await db.execute(text("""
            SELECT t.id, t.name, t.name_jp, t.vendor,
                   COALESCE(tc.name_jp, '') as category_name_jp
            FROM tools t
            LEFT JOIN tool_categories tc ON t.category_id = tc.id
            WHERE t.is_active = true
            ORDER BY t.name_jp
        """))
        rows = result.fetchall()
    except Exception as e:
        logger.error("Failed to query tools with category join: %s", e)
        # Rollback the failed transaction before retrying
        await db.rollback()
        try:
            result = await db.execute(text(
                "SELECT id, name, name_jp, vendor, '' FROM tools WHERE is_active = true ORDER BY name_jp"
            ))
            rows = result.fetchall()
        except Exception as e2:
            logger.error("Fallback tool query also failed: %s", e2)
            await db.rollback()
            raise HTTPException(500, f"ツール一覧の取得に失敗しました: {e2}")

    return [
        ToolListItem(
            id=row[0],
            name=row[1] or "",
            name_jp=row[2] or "",
            vendor=row[3] or "",
            category_name_jp=row[4] or "",
        )
        for row in rows
    ]
